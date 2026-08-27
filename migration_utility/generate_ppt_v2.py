"""
Generate Reference Architecture PowerPoint — v2
SQL Server → Databricks Unity Catalog Migration Studio
Full-featured: 14 UI tabs, 85+ API endpoints, Self-Healing, CDC/SDP, etc.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colours ─────────────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xFD, 0xFD, 0xFD)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
AZURE_BLUE  = RGBColor(0x00, 0x78, 0xD4)
DBR_RED     = RGBColor(0xFF, 0x3E, 0x21)
LANDING_CLR = RGBColor(0x60, 0x7D, 0x8B)
BRONZE_CLR  = RGBColor(0xCD, 0x7F, 0x32)
SILVER_CLR  = RGBColor(0xAA, 0xA9, 0xAD)
GOLD_CLR    = RGBColor(0xFF, 0xD7, 0x00)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
PURPLE      = RGBColor(0x7B, 0x1F, 0xA2)
TEAL        = RGBColor(0x00, 0x96, 0x88)
ORANGE      = RGBColor(0xE6, 0x51, 0x00)
RED         = RGBColor(0xC6, 0x28, 0x28)
INFRA_BLUE  = RGBColor(0x15, 0x65, 0xC0)
DEEP_BLUE   = RGBColor(0x0D, 0x47, 0xA1)
AMBER       = RGBColor(0xFF, 0x8F, 0x00)
INDIGO      = RGBColor(0x30, 0x3F, 0x9F)
PINK        = RGBColor(0xAD, 0x14, 0x57)
CYAN        = RGBColor(0x00, 0x83, 0x8F)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────

def box(slide, l, t, w, h, fill, text="", sz=9, fc=WHITE, bold=False,
        align=PP_ALIGN.CENTER, line_clr=None, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_clr:
        s.line.color.rgb = line_clr; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = radius
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(2)
    if text:
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.color.rgb = fc; r.font.bold = bold
    return s


def mbox(slide, l, t, w, h, lines, sz=8, fc=DARK_GRAY, fill=None, line_clr=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_clr:
        s.line.color.rgb = line_clr; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = 0.04
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        r = p.add_run(); r.text = ln; r.font.size = Pt(sz); r.font.color.rgb = fc
    return s


def txt(slide, l, t, w, h, text, sz=9, fc=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(sz); r.font.color.rgb = fc; r.font.bold = bold
    return tb


def arrow_r(slide, x1, y, x2, color=MED_GRAY):
    ln = slide.shapes.add_connector(1, x1, y, x2 - Inches(0.12), y)
    ln.line.color.rgb = color; ln.line.width = Pt(2)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  x2 - Inches(0.14), y - Inches(0.06), Inches(0.14), Inches(0.12))
    tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = color; tri.line.fill.background()


def arrow_d(slide, x, y1, y2, color=MED_GRAY):
    ln = slide.shapes.add_connector(1, x, y1, x, y2 - Inches(0.12))
    ln.line.color.rgb = color; ln.line.width = Pt(1.5)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  x - Inches(0.06), y2 - Inches(0.12), Inches(0.12), Inches(0.12))
    tri.rotation = 180; tri.fill.solid(); tri.fill.fore_color.rgb = color; tri.line.fill.background()


def section_header(slide, x, y, w, title, color):
    box(slide, x, y, w, Inches(0.28), color, title, sz=9, fc=WHITE, bold=True)


def feature_card(slide, x, y, w, h, title, items, color, light):
    box(slide, x, y, w, Inches(0.26), color, title, sz=8, fc=WHITE, bold=True)
    mbox(slide, x, y + Inches(0.28), w, h - Inches(0.28), items,
         sz=6.5, fill=light, line_clr=color)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Reference Architecture Overview
# ═══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6])

# ── Title Bar ───────────────────────────────────────────────────────────────
box(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.55), DEEP_BLUE,
    "Reference Architecture  —  SQL Server → Databricks Unity Catalog Migration Studio  |  v2.0",
    sz=15, fc=WHITE, bold=True)

# Subtitle bar
box(sl, Inches(0), Inches(0.55), Inches(13.333), Inches(0.2), RGBColor(0xE3, 0xF2, 0xFD),
    "14 UI Pages  •  85+ API Endpoints  •  Metadata-Driven Medallion  •  CDC/SDP  •  Self-Healing  •  Auto Infrastructure  •  Data Quality  •  Reconciliation",
    sz=7, fc=DEEP_BLUE, bold=True)


# ═════════════════════════════════════════════════════════════════════════════
#  ROW 1 — Source → Flask AI Workflow Manager → Databricks (Medallion)
# ═════════════════════════════════════════════════════════════════════════════

row1_y = Inches(0.9)

# ── Source: Azure SQL ────────────────────────────────────────────────────────
src_x = Inches(0.15)
box(sl, src_x, row1_y, Inches(1.5), Inches(0.3), AZURE_BLUE,
    "Azure SQL / SQL Server", sz=8, fc=WHITE, bold=True)
mbox(sl, src_x, row1_y + Inches(0.33), Inches(1.5), Inches(0.7),
     ["4 Source Types:", "• Azure SQL Database", "• SQL Server (on-prem)",
      "• SQL Managed Instance", "• Synapse Analytics"],
     sz=6.5, fill=RGBColor(0xE3, 0xF2, 0xFD), line_clr=AZURE_BLUE)

# Stored Procs / Views / UDFs box under source
box(sl, src_x, row1_y + Inches(1.1), Inches(1.5), Inches(0.22), ORANGE,
    "SPs  •  Views  •  UDFs", sz=7, fc=WHITE, bold=True)
mbox(sl, src_x, row1_y + Inches(1.35), Inches(1.5), Inches(0.35),
     ["AI SP Converter", "OpenAI / Gemini / Claude", "SQL → PySpark notebooks"],
     sz=6, fill=RGBColor(0xFF, 0xF3, 0xE0), line_clr=ORANGE)

# Arrow → Manager
arrow_r(sl, Inches(1.7), row1_y + Inches(0.55), Inches(1.95), AZURE_BLUE)

# ── Flask AI Workflow Manager ────────────────────────────────────────────────
mgr_x = Inches(1.95)
mgr_w = Inches(2.0)
box(sl, mgr_x, row1_y, mgr_w, Inches(0.3), PURPLE,
    "Migration Studio (Flask)", sz=8, fc=WHITE, bold=True)
mbox(sl, mgr_x, row1_y + Inches(0.33), mgr_w, Inches(1.37),
     ["AI Workflow Manager (14 pages)",
      "─────────────────────────",
      "• Pipeline Studio + Job Manager",
      "• MetadataFlow orchestration",
      "• Auto-init from deployconfig",
      "• Notebook gen & deploy",
      "• Databricks job submit + poll",
      "• CDC / SDP mode config",
      "• Self-Healing Bot (auto-fix)",
      "• Progress Tracker (7 steps)",
      "• Reports & Email alerts"],
     sz=6.5, fill=RGBColor(0xF3, 0xE5, 0xF5), line_clr=PURPLE)

# Arrow → Databricks
arrow_r(sl, Inches(4.0), row1_y + Inches(0.55), Inches(4.35), PURPLE)
txt(sl, Inches(4.0), row1_y + Inches(0.2), Inches(0.5), Inches(0.2),
    "REST API", sz=6, fc=MED_GRAY, align=PP_ALIGN.CENTER)

# ── Databricks Container ────────────────────────────────────────────────────
dbr_x = Inches(4.35)
dbr_w = Inches(8.8)
dbr_container = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     dbr_x, row1_y - Inches(0.05), dbr_w, Inches(2.75))
dbr_container.fill.solid()
dbr_container.fill.fore_color.rgb = NEAR_WHITE
dbr_container.line.color.rgb = DBR_RED
dbr_container.line.width = Pt(2)
dbr_container.adjustments[0] = 0.02
txt(sl, dbr_x + Inches(0.1), row1_y - Inches(0.02), Inches(3.5), Inches(0.22),
    "Databricks Unity Catalog  —  Multi-Catalog Architecture", sz=9, fc=DBR_RED, bold=True)

# ── Medallion Layers ────────────────────────────────────────────────────────
layer_y = row1_y + Inches(0.25)
layer_w = Inches(1.85)
gap = Inches(0.18)

layers = [
    ("dev_volumes",   "UC Volumes (Landing)",  LANDING_CLR, RGBColor(0xEC, 0xEF, 0xF1),
     ["/Volumes/dev_volumes/", "  hr/landing/{table}", "Parquet staging"]),
    ("bronze",        "Bronze (Raw + Audit)",   BRONZE_CLR,  RGBColor(0xFB, 0xE9, 0xD0),
     ["bronze.hr.{table}", "Delta Tables", "+ audit cols + DQ"]),
    ("silver",        "Silver (Cleansed)",      SILVER_CLR,  RGBColor(0xEC, 0xEF, 0xF1),
     ["silver.hr.{table}", "Dedup + null handling", "7-check DQ scoring"]),
    ("gold (future)", "Gold (Business KPIs)",   GOLD_CLR,    RGBColor(0xFF, 0xF8, 0xE1),
     ["Aggregated datasets", "BI-ready views", "Power BI / Tableau"]),
]

lx = dbr_x + Inches(0.12)
for i, (cat, subtitle, clr, light, items) in enumerate(layers):
    x = lx + i * (layer_w + gap)
    fc_title = WHITE if cat != "gold (future)" else DARK_GRAY
    box(sl, x, layer_y, layer_w, Inches(0.24), clr, cat, sz=8, fc=fc_title, bold=True)
    box(sl, x, layer_y + Inches(0.26), layer_w, Inches(0.18), light, subtitle,
        sz=6.5, fc=DARK_GRAY, bold=True, line_clr=clr)
    mbox(sl, x, layer_y + Inches(0.46), layer_w, Inches(0.42), items,
         sz=6.5, fill=light, line_clr=clr)
    if i < len(layers) - 1:
        arrow_r(sl, x + layer_w, layer_y + Inches(0.55), x + layer_w + gap, clr)

# ── Metadata Tables (inside Databricks) ─────────────────────────────────────
meta_y = row1_y + Inches(1.2)
box(sl, dbr_x + Inches(0.12), meta_y, dbr_w - Inches(0.24), Inches(0.22),
    GREEN, "Metadata Delta Tables  •  Reconciliation  •  Logging  •  DQ Metrics", sz=7.5, fc=WHITE, bold=True)

meta_tables = ["wf_pipeline_metadata", "wf_job_metadata", "wf_run_history",
               "wf_watermark_metadata", "wf_source_tables",
               "Reconciliation.*", "ExecutionLog.*"]
mt_w = (dbr_w - Inches(0.24) - Inches(0.07) * 6) / 7
for i, tbl in enumerate(meta_tables):
    mx = dbr_x + Inches(0.12) + i * (mt_w + Inches(0.07))
    clr = GREEN if i < 5 else ORANGE
    box(sl, mx, meta_y + Inches(0.24), mt_w, Inches(0.26),
        RGBColor(0xE8, 0xF5, 0xE9) if i < 5 else RGBColor(0xFF, 0xF3, 0xE0),
        tbl, sz=5.5, fc=DARK_GRAY, line_clr=clr)

# ── Notebooks (inside Databricks) ───────────────────────────────────────────
nb_y = meta_y + Inches(0.55)
box(sl, dbr_x + Inches(0.12), nb_y, dbr_w - Inches(0.24), Inches(0.22),
    TEAL, "Auto-Generated Notebooks  (Standard PySpark  +  SDP Mode  +  CDC)", sz=7.5, fc=WHITE, bold=True)

notebooks = [
    ("00_Orchestrator",    "Chains stages\nDLT or Spark"),
    ("01_Extract",         "JDBC + Watermark\nParquet → Volumes"),
    ("02_Bronze/SDP",      "Landing → Bronze\nAudit + DQ checks"),
    ("03_Silver",          "Bronze → Silver\n7-check DQ score"),
    ("04_Reconciliation",  "Source vs Bronze\nSUM per col"),
    ("05_ExecutionLog",    "Audit trail\nPer-job results"),
]
nb_w = (dbr_w - Inches(0.24) - Inches(0.1) * 5) / 6
for i, (name, desc) in enumerate(notebooks):
    nx = dbr_x + Inches(0.12) + i * (nb_w + Inches(0.1))
    box(sl, nx, nb_y + Inches(0.24), nb_w, Inches(0.2),
        RGBColor(0xE0, 0xF2, 0xF1), name, sz=6, fc=DARK_GRAY, bold=True, line_clr=TEAL)
    mbox(sl, nx, nb_y + Inches(0.46), nb_w, Inches(0.34), desc.split("\n"),
         sz=6, fill=RGBColor(0xE0, 0xF2, 0xF1), line_clr=TEAL)
    if i < len(notebooks) - 1:
        arrow_r(sl, nx + nb_w, nb_y + Inches(0.5),
                nx + nb_w + Inches(0.1), TEAL)


# ═════════════════════════════════════════════════════════════════════════════
#  ROW 2 — Left: Deploy Infra | Middle: UI Pages (14 tabs) | Right: 4 Pillars
# ═════════════════════════════════════════════════════════════════════════════

row2_y = Inches(3.85)

# ── Deploy Infrastructure ────────────────────────────────────────────────────
box(sl, Inches(0.15), row2_y, Inches(1.95), Inches(0.26), INFRA_BLUE,
    "Deploy Infrastructure", sz=8, fc=WHITE, bold=True)
mbox(sl, Inches(0.15), row2_y + Inches(0.28), Inches(1.95), Inches(0.78),
     ["AutoInfraCreation (Azure SDK)",
      "1. ADLS Gen2 + HNS",
      "2. Access Connector (MI)",
      "3. RBAC role assignment",
      "4. UC Storage Credential",
      "5. External Locations",
      "6. Catalogs + MANAGED LOC",
      "7. UC Volumes (landing)"],
     sz=6, fill=RGBColor(0xE3, 0xF2, 0xFD), line_clr=INFRA_BLUE)

# ── Self-Healing Bot ─────────────────────────────────────────────────────────
box(sl, Inches(0.15), row2_y + Inches(1.15), Inches(0.93), Inches(0.22), RED,
    "Self-Healing Bot", sz=6.5, fc=WHITE, bold=True)
mbox(sl, Inches(0.15), row2_y + Inches(1.39), Inches(0.93), Inches(0.45),
     ["Error classification", "Auto-retry + backoff",
      "Cluster restart", "Restore points"],
     sz=5.5, fill=RGBColor(0xFF, 0xEB, 0xEE), line_clr=RED)

# ── CDC / SDP Config ─────────────────────────────────────────────────────────
box(sl, Inches(1.13), row2_y + Inches(1.15), Inches(0.97), Inches(0.22), TEAL,
    "CDC / SDP Engine", sz=6.5, fc=WHITE, bold=True)
mbox(sl, Inches(1.13), row2_y + Inches(1.39), Inches(0.97), Inches(0.45),
     ["Change Tracking CDC", "SDP Auto Loader",
      "Watermark-based incr.", "Primary key config"],
     sz=5.5, fill=RGBColor(0xE0, 0xF2, 0xF1), line_clr=TEAL)

# ── UI Pages / 14 Tabs ──────────────────────────────────────────────────────
ui_x = Inches(2.35)
ui_w = Inches(4.8)
box(sl, ui_x, row2_y, ui_w, Inches(0.26), PURPLE,
    "Migration Studio UI  —  14 Feature-Rich Pages", sz=8, fc=WHITE, bold=True)

# Column 1 of tabs
tab_col1 = [
    ("📊", "Dashboard",         "Pipeline health, stats, metrics"),
    ("🔄", "MetadataFlow",      "Init metadata tables, sync"),
    ("⚡", "Pipeline Studio",    "Create & run pipelines, bulk ops"),
    ("📋", "Job Manager",       "CRUD jobs, rerun, watermarks"),
    ("📈", "Reports",           "Email reports, job summaries"),
    ("📍", "Progress Tracker",  "7-step migration progress"),
    ("🔐", "Audit & Compliance","Event log, severity filters"),
]
tab_col2 = [
    ("✨", "Data Quality",      "DQ rules, expectations, scores"),
    ("🔀", "Schema Comparison", "Drift detection, type mismatch"),
    ("📐", "Reconciliation",    "SUM validation, source vs target"),
    ("⚙️", "Settings",          "Creds, source DB, metadata loc"),
    ("🔧", "Convert to PySpark","SP/View/UDF → PySpark (AI)"),
    ("🚀", "Deploy Notebooks",  "Batch upload to workspace"),
    ("📦", "Data Migration",    "JDBC → Parquet → COPY INTO"),
]

tab_h = Inches(0.22)
tab_w1 = Inches(2.32)
tab_w2 = Inches(2.32)
for i, (icon, name, desc) in enumerate(tab_col1):
    ty = row2_y + Inches(0.3) + i * (tab_h + Inches(0.02))
    bg = RGBColor(0xF3, 0xE5, 0xF5) if i % 2 == 0 else WHITE
    box(sl, ui_x, ty, tab_w1, tab_h, bg,
        f"{icon} {name}", sz=6.5, fc=DARK_GRAY, bold=True, align=PP_ALIGN.LEFT,
        line_clr=RGBColor(0xE0, 0xE0, 0xE0))

for i, (icon, name, desc) in enumerate(tab_col2):
    ty = row2_y + Inches(0.3) + i * (tab_h + Inches(0.02))
    bg = RGBColor(0xF3, 0xE5, 0xF5) if i % 2 == 0 else WHITE
    box(sl, ui_x + tab_w1 + Inches(0.08), ty, tab_w2, tab_h, bg,
        f"{icon} {name}", sz=6.5, fc=DARK_GRAY, bold=True, align=PP_ALIGN.LEFT,
        line_clr=RGBColor(0xE0, 0xE0, 0xE0))

# ── 4 Pillars ───────────────────────────────────────────────────────────────
pill_x0 = Inches(7.4)
pill_w = Inches(1.42)
pill_gap = Inches(0.1)

pillars = [
    ("Security", RED, RGBColor(0xFF, 0xEB, 0xEE),
     ["PAT token auth", "Encrypted JDBC",
      "UC RBAC per catalog", "MANAGED LOCATION"]),
    ("Governance", GREEN, RGBColor(0xE8, 0xF5, 0xE9),
     ["Reconciliation catalog", "Execution logging",
      "Delta versioning", "Audit trail"]),
    ("Reliability", PURPLE, RGBColor(0xF3, 0xE5, 0xF5),
     ["Self-healing bot", "Per-stage rerun",
      "DQ scoring (12+)", "Restore points"]),
    ("Scale", TEAL, RGBColor(0xE0, 0xF2, 0xF1),
     ["Metadata-driven", "SDP + PySpark",
      "Multi-catalog", "CDC watermarks"]),
]

for i, (title, clr, light, items) in enumerate(pillars):
    px = pill_x0 + i * (pill_w + pill_gap)
    box(sl, px, row2_y, pill_w, Inches(0.26), clr, title, sz=7.5, fc=WHITE, bold=True)
    mbox(sl, px, row2_y + Inches(0.28), pill_w, Inches(0.62), items,
         sz=6, fill=light, line_clr=clr)

# ── Stats Ribbon ─────────────────────────────────────────────────────────────
stats_y = row2_y + Inches(0.97)
stats = [
    ("85+", "API Endpoints", AZURE_BLUE),
    ("14", "UI Pages", PURPLE),
    ("12", "Python Modules", TEAL),
    ("6", "Notebooks", GREEN),
    ("10+", "Healing Actions", RED),
    ("12+", "DQ Checks", ORANGE),
]
stat_w = Inches(0.93)
for i, (num, label, clr) in enumerate(stats):
    sx = pill_x0 + i * (stat_w + Inches(0.08))
    box(sl, sx, stats_y, stat_w, Inches(0.24), clr, f"{num}  {label}",
        sz=6.5, fc=WHITE, bold=True)


# ═════════════════════════════════════════════════════════════════════════════
#  ROW 3 — Tech Stack + Data Flow + Key Configs
# ═════════════════════════════════════════════════════════════════════════════

row3_y = Inches(5.7)

# ── Tech Stack ──────────────────────────────────────────────────────────────
box(sl, Inches(0.15), row3_y, Inches(4.0), Inches(0.22), DARK_GRAY,
    "Technology Stack", sz=8, fc=WHITE, bold=True)

tech_items = [
    ("Azure SQL / SQL Server", "4 source types (Azure SQL, SQL Server, MI, Synapse)"),
    ("Python 3.11 / Flask",    "Web UI + REST API (85+ endpoints)"),
    ("Databricks REST API",    "Job submit, poll, notebook deploy, cluster ops"),
    ("UC Volumes + Delta",     "Landing (Parquet) → Bronze/Silver (Delta)"),
    ("Lakeflow Spark Declarative Pipelines",      "SDP Auto Loader + expectations + CDC"),
    ("deployconfig.json",      "Centralised config — auto-loads on startup"),
]

row_h = Inches(0.2)
for i, (tech, desc) in enumerate(tech_items):
    ry = row3_y + Inches(0.24) + i * row_h
    bg = RGBColor(0xE3, 0xF2, 0xFD) if i % 2 == 0 else WHITE
    box(sl, Inches(0.15), ry, Inches(1.65), row_h, bg, tech,
        sz=6, fc=DARK_GRAY, bold=True, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    box(sl, Inches(1.8), ry, Inches(2.35), row_h, bg, desc,
        sz=6, fc=DARK_GRAY, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))

# ── End-to-End Data Flow ────────────────────────────────────────────────────
flow_x = Inches(4.4)
box(sl, flow_x, row3_y, Inches(5.5), Inches(0.22), AZURE_BLUE,
    "End-to-End Data Flow", sz=8, fc=WHITE, bold=True)

flow_steps = [
    ("Azure\nSQL", AZURE_BLUE),
    ("JDBC\nExtract", PURPLE),
    ("Volumes\nLanding", LANDING_CLR),
    ("Bronze\nDelta", BRONZE_CLR),
    ("DQ +\nRecon", ORANGE),
    ("Silver\nDelta", SILVER_CLR),
    ("Exec\nLog", GREEN),
]

fs_w = Inches(0.68)
fs_gap = Inches(0.07)
fs_y = row3_y + Inches(0.28)
for i, (label, clr) in enumerate(flow_steps):
    fx = flow_x + Inches(0.1) + i * (fs_w + fs_gap)
    box(sl, fx, fs_y, fs_w, Inches(0.46), clr, label, sz=6.5, fc=WHITE, bold=True)
    if i < len(flow_steps) - 1:
        arrow_r(sl, fx + fs_w, fs_y + Inches(0.23), fx + fs_w + fs_gap, clr)

# Pipeline modes annotation
mbox(sl, flow_x + Inches(0.1), fs_y + Inches(0.52), Inches(5.3), Inches(0.32),
     ["Standard Mode: Extract → Bronze → Silver (3 stages)  |  SDP Mode: Extract → SDP Bronze+Silver (2 stages)",
      "CDC: Change Tracking + Watermark-based incremental  |  Rerun: Per-job rerun from first failure"],
     sz=6, fill=RGBColor(0xE3, 0xF2, 0xFD), line_clr=AZURE_BLUE)

# ── Config / Key Integration Points ─────────────────────────────────────────
cfg_x = Inches(10.15)
cfg_w = Inches(3.0)
box(sl, cfg_x, row3_y, cfg_w, Inches(0.22), DARK_GRAY,
    "Key Integration Points", sz=8, fc=WHITE, bold=True)

configs = [
    "deployconfig.json — all creds + catalog map",
    "Auto-init on page load (MetadataFlow)",
    "Settings page: Source DB + Metadata loc",
    "Schema Comparison: drift detection",
    "Data Quality: 12+ checks + __dq_metrics",
    "Reconciliation: SUM validation per col",
]
for i, item in enumerate(configs):
    ry = row3_y + Inches(0.24) + i * Inches(0.19)
    box(sl, cfg_x, ry, cfg_w, Inches(0.18),
        WHITE if i % 2 == 0 else LIGHT_GRAY, "• " + item,
        sz=5.5, fc=DARK_GRAY, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))

# ── Footer ──────────────────────────────────────────────────────────────────
txt(sl, Inches(0.2), Inches(7.2), Inches(13), Inches(0.2),
    "SQL → Databricks Migration Studio v2.0  |  Metadata-Driven Medallion  |  CDC + SDP  |  Self-Healing  |  14 UI Pages  |  85+ Endpoints  |  deployconfig.json",
    sz=7, fc=MED_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — What's New Since v1 (Comprehensive Changelog)
# ═══════════════════════════════════════════════════════════════════════════════

sl2 = prs.slides.add_slide(prs.slide_layouts[6])

box(sl2, Inches(0), Inches(0), Inches(13.333), Inches(0.5), GREEN,
    "What's New in v2.0  —  All Features Added Since Initial Architecture (v1)",
    sz=15, fc=WHITE, bold=True)

# ═══════════════ COLUMN 1 — Architecture & Backend ══════════════════════════

col1_x = Inches(0.15)
col1_w = Inches(4.2)
cy = Inches(0.6)

section_header(sl2, col1_x, cy, col1_w, "🏗️ Architecture & Backend", DEEP_BLUE)
cy += Inches(0.34)

feature_card(sl2, col1_x, cy, col1_w, Inches(0.92),
    "Settings Page — Centralized Configuration",
    ["Source Database Connection (server, db, user, pass)",
     "Test Connection button with server version display",
     "Metadata Tables Location (catalog + schema)",
     "Auto-save to deployconfig.json on change",
     "All pages auto-populate from saved config"],
    INFRA_BLUE, RGBColor(0xE3, 0xF2, 0xFD))
cy += Inches(0.97)

feature_card(sl2, col1_x, cy, col1_w, Inches(0.92),
    "Auto-Init from deployconfig.json",
    ["Backend _restore_from_deploy_config() on module load",
     "Auto-detects SQL Warehouse (inline, no NameError)",
     "Sets _metadata_initialized=True if config exists",
     "Frontend _autoInitFromConfig() on page load",
     "No need to visit MetadataFlow tab first"],
    PURPLE, RGBColor(0xF3, 0xE5, 0xF5))
cy += Inches(0.97)

feature_card(sl2, col1_x, cy, col1_w, Inches(0.92),
    "Pipeline Studio Redesign",
    ["Removed Connect Data Source panel (uses Settings)",
     "Two-column table picker: Available ↔ Selected",
     "Blue checkboxes, hover effects, search + filter",
     "Compact Cluster bar (was 50/50 → single row)",
     "Auto-populate tables on Discover click"],
    TEAL, RGBColor(0xE0, 0xF2, 0xF1))
cy += Inches(0.97)

feature_card(sl2, col1_x, cy, col1_w, Inches(0.92),
    "Job Manager — Auto Run on Databricks",
    ["Auto-detect deployment from /api/workflow/notebooks/status",
     "Fallback credential chain: UI → Settings → deployconfig",
     "Auto-detect running clusters for execution",
     "Background poller updates JOB_REGISTRY status",
     "10s auto-refresh when running jobs exist"],
    AMBER, RGBColor(0xFF, 0xF8, 0xE1))
cy += Inches(0.97)

feature_card(sl2, col1_x, cy, col1_w, Inches(0.72),
    "SQL Objects Panel Redesign",
    ["Rounded card layout with icon + count pill",
     "Filter pills (All, None, SP, Views, UDFs)",
     "Larger checkboxes, type descriptions, tooltips",
     "Grouped sections with collapse/expand"],
    ORANGE, RGBColor(0xFF, 0xF3, 0xE0))


# ═══════════════ COLUMN 2 — Pipeline & Data Features ═══════════════════════

col2_x = Inches(4.55)
col2_w = Inches(4.2)
cy2 = Inches(0.6)

section_header(sl2, col2_x, cy2, col2_w, "⚡ Pipeline & Data Features", PURPLE)
cy2 += Inches(0.34)

feature_card(sl2, col2_x, cy2, col2_w, Inches(0.78),
    "CDC + SDP Engine",
    ["Change Tracking CDC mode for incremental loads",
     "SDP Auto Loader with Delta expectations",
     "Watermark-based incremental with auto-advance",
     "Primary key configuration in deployconfig"],
    TEAL, RGBColor(0xE0, 0xF2, 0xF1))
cy2 += Inches(0.83)

feature_card(sl2, col2_x, cy2, col2_w, Inches(0.78),
    "Data Quality Framework (12+ Checks)",
    ["Bronze: empty file, null-key, dups, schema drift, quarantine",
     "Silver: null removal, high-null cols, dedup, trim, anomaly",
     "DQ Score % → saved to __dq_metrics Delta table",
     "Per-row: __dq_status (passed/warn) + __is_quarantined"],
    GREEN, RGBColor(0xE8, 0xF5, 0xE9))
cy2 += Inches(0.83)

feature_card(sl2, col2_x, cy2, col2_w, Inches(0.65),
    "Self-Healing Bot — Auto Recovery",
    ["Error classification: OOM, Auth, Schema, Timeout",
     "Healing: retry (exp backoff), restart, scale up, rollback",
     "Restore points: named snapshots for safe recovery"],
    RED, RGBColor(0xFF, 0xEB, 0xEE))
cy2 += Inches(0.70)

feature_card(sl2, col2_x, cy2, col2_w, Inches(0.65),
    "Schema Comparison & Drift Detection",
    ["Source vs target column comparison",
     "Type mismatch + nullable change detection",
     "Smart recommendations per column"],
    INDIGO, RGBColor(0xE8, 0xEA, 0xF6))
cy2 += Inches(0.70)

feature_card(sl2, col2_x, cy2, col2_w, Inches(0.65),
    "Reconciliation & Execution Logging",
    ["04_Reconciliation: SUM per numeric col validation",
     "05_ExecutionLog: per-job audit trail",
     "Saves to dedicated catalogs (configurable)"],
    ORANGE, RGBColor(0xFF, 0xF3, 0xE0))
cy2 += Inches(0.70)

feature_card(sl2, col2_x, cy2, col2_w, Inches(0.52),
    "Deploy Infrastructure (AutoInfraCreation)",
    ["7-step Azure + UC setup via Azure Python SDK",
     "ADLS Gen2, Access Connector, RBAC, Catalogs, Volumes"],
    INFRA_BLUE, RGBColor(0xE3, 0xF2, 0xFD))


# ═══════════════ COLUMN 3 — Bug Fixes + Before/After ══════════════════════

col3_x = Inches(8.95)
col3_w = Inches(4.2)
cy3 = Inches(0.6)

section_header(sl3 := sl2, col3_x, cy3, col3_w, "🔧 Fixes & Improvements", RED)
cy3 += Inches(0.34)

feature_card(sl2, col3_x, cy3, col3_w, Inches(0.65),
    "Metadata Table Location Fix",
    ["Fixed TABLE_OR_VIEW_NOT_FOUND for wf_pipeline_metadata",
     "Tables exist in admin_source.configtables (not main.default)",
     "Added metadata_catalog / metadata_schema to deployconfig"],
    MED_GRAY, LIGHT_GRAY)
cy3 += Inches(0.70)

feature_card(sl2, col3_x, cy3, col3_w, Inches(0.52),
    "Job Status Polling Fix",
    ["Poller now updates JOB_REGISTRY (was updating only JOB_RUNS)",
     "Auto-refresh timer (10s) when running jobs exist"],
    MED_GRAY, LIGHT_GRAY)
cy3 += Inches(0.57)

feature_card(sl2, col3_x, cy3, col3_w, Inches(0.52),
    "NameError Fix — Backend Auto-Restore",
    ["_restore_from_deploy_config() inlined warehouse detection",
     "No longer calls _dbr_session() before it's defined"],
    MED_GRAY, LIGHT_GRAY)
cy3 += Inches(0.57)

feature_card(sl2, col3_x, cy3, col3_w, Inches(0.52),
    "JDBC ORDER BY Fix",
    ["Removed ORDER BY ORDINAL_POSITION from JDBC subquery",
     "SQL Server forbids ORDER BY in derived tables"],
    MED_GRAY, LIGHT_GRAY)
cy3 += Inches(0.57)

# ── Before vs After ─────────────────────────────────────────────────────────
box(sl2, col3_x, cy3, col3_w, Inches(0.26), DEEP_BLUE,
    "Before (v1) → After (v2)", sz=8, fc=WHITE, bold=True)
cy3 += Inches(0.28)

before_after = [
    ("UI Pages",      "Basic",          "14 feature-rich tabs"),
    ("API Endpoints",  "~30",           "85+ routes"),
    ("Config",         "Hardcoded",     "deployconfig.json auto-init"),
    ("MetadataFlow",   "Manual init",    "Auto-init on page load"),
    ("Pipeline UX",    "Dropdown picker","Two-panel table selector"),
    ("Cluster UI",     "50% width card", "Compact single-row bar"),
    ("CDC/SDP",        "Basic watermark","CDC + SDP + Auto Loader"),
    ("DQ Checks",      "Basic",         "12+ checks + DQ Score"),
    ("Self-Healing",   "None",          "10+ auto-recovery actions"),
    ("Schema Compare", "None",          "Drift detection + smart recs"),
    ("Settings",       "None",          "Source DB + Metadata config"),
]

# Header row
box(sl2, col3_x, cy3, Inches(1.0), Inches(0.18),
    RGBColor(0xE0, 0xE0, 0xE0), "Feature", sz=6, fc=DARK_GRAY, bold=True,
    align=PP_ALIGN.LEFT, line_clr=RGBColor(0xBD, 0xBD, 0xBD))
box(sl2, col3_x + Inches(1.0), cy3, Inches(1.3), Inches(0.18),
    RGBColor(0xFF, 0xEB, 0xEE), "v1 (Before)", sz=6, fc=RED, bold=True,
    align=PP_ALIGN.LEFT, line_clr=RGBColor(0xBD, 0xBD, 0xBD))
box(sl2, col3_x + Inches(2.3), cy3, Inches(1.9), Inches(0.18),
    RGBColor(0xE8, 0xF5, 0xE9), "v2 (After)", sz=6, fc=GREEN, bold=True,
    align=PP_ALIGN.LEFT, line_clr=RGBColor(0xBD, 0xBD, 0xBD))
cy3 += Inches(0.18)

for j, (feat, before, after) in enumerate(before_after):
    bg = WHITE if j % 2 == 0 else RGBColor(0xFA, 0xFA, 0xFA)
    box(sl2, col3_x, cy3, Inches(1.0), Inches(0.16), bg, feat,
        sz=5.5, fc=DARK_GRAY, bold=True, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    box(sl2, col3_x + Inches(1.0), cy3, Inches(1.3), Inches(0.16), bg, before,
        sz=5.5, fc=MED_GRAY, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    box(sl2, col3_x + Inches(2.3), cy3, Inches(1.9), Inches(0.16), bg, after,
        sz=5.5, fc=GREEN, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    cy3 += Inches(0.16)

# Footer
txt(sl2, Inches(0.2), Inches(7.2), Inches(13), Inches(0.2),
    "SQL → Databricks Migration Studio v2.0  |  Complete feature changelog from v1 → v2",
    sz=7, fc=MED_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — Detailed UI Pages Map
# ═══════════════════════════════════════════════════════════════════════════════

sl3 = prs.slides.add_slide(prs.slide_layouts[6])

box(sl3, Inches(0), Inches(0), Inches(13.333), Inches(0.5), PURPLE,
    "Migration Studio  —  14 UI Pages  •  Feature Map",
    sz=15, fc=WHITE, bold=True)

# Page cards — 3 rows x 5 cols
pages = [
    # Row 1
    ("📊 Dashboard", AZURE_BLUE, RGBColor(0xE3, 0xF2, 0xFD),
     ["Pipeline health overview", "Job status + exec metrics", "Quick-access navigation",
      "Total pipelines / jobs / runs"]),
    ("🔄 MetadataFlow", PURPLE, RGBColor(0xF3, 0xE5, 0xF5),
     ["Init metadata Delta tables", "Load / Sync to Databricks", "Auto-init from config",
      "5 metadata tables managed"]),
    ("⚡ Pipeline Studio", TEAL, RGBColor(0xE0, 0xF2, 0xF1),
     ["Medallion architecture viz", "Two-panel table selector", "Quick Create (bulk)",
      "Run on Databricks + rerun"]),
    ("📋 Job Manager", AMBER, RGBColor(0xFF, 0xF8, 0xE1),
     ["Create / update / delete jobs", "Watermark management", "Per-job rerun on fail",
      "10s auto-refresh polling"]),
    ("🔧 Convert to PySpark", ORANGE, RGBColor(0xFF, 0xF3, 0xE0),
     ["AI: OpenAI / Gemini / Claude", "SP + View + UDF support", "Multi-file export",
      "Syntax-highlighted preview"]),
    # Row 2
    ("🚀 Deploy Notebooks", INFRA_BLUE, RGBColor(0xE3, 0xF2, 0xFD),
     ["Batch upload to workspace", "Path auto-create", "Per-notebook status",
      "Direct workspace URLs"]),
    ("📦 Unity Catalog", DEEP_BLUE, RGBColor(0xE3, 0xF2, 0xFD),
     ["Browse catalogs + schemas", "Table preview (sample rows)", "Custom SQL execution",
      "Warehouse auto-filter"]),
    ("📤 Data Migration", BRONZE_CLR, RGBColor(0xFB, 0xE9, 0xD0),
     ["JDBC → Parquet → COPY INTO", "Full + Incremental load", "Parallel workers (1-8)",
      "Per-table progress bars"]),
    ("🏗️ Medallion Arch", LANDING_CLR, RGBColor(0xEC, 0xEF, 0xF1),
     ["Generate 6 notebooks", "Standard + SDP mode", "Multi-catalog support",
      "Run Orchestrator as job"]),
    ("🤖 Self-Healing Bot", RED, RGBColor(0xFF, 0xEB, 0xEE),
     ["Error classification (8 types)", "Auto-retry + backoff", "Restore points",
      "Rules engine + custom rules"]),
    # Row 3
    ("📈 Reports", GREEN, RGBColor(0xE8, 0xF5, 0xE9),
     ["Job execution reports", "Email notifications", "Duration + row metrics",
      "Success rate tracking"]),
    ("📍 Progress Tracker", INDIGO, RGBColor(0xE8, 0xEA, 0xF6),
     ["7-step migration workflow", "Step indicators + status", "Current step highlight",
      "End-to-end visibility"]),
    ("🔐 Audit & Compliance", PINK, RGBColor(0xFC, 0xE4, 0xEC),
     ["Event log with timestamps", "Severity: INFO→CRITICAL", "Filter by time + type",
      "Compliance dashboard"]),
    ("✨ Data Quality", GREEN, RGBColor(0xE8, 0xF5, 0xE9),
     ["DQ rules library", "12+ quality checks", "DQ Score + __dq_metrics",
      "Quality trend charts"]),
    ("🔀 Schema Compare", CYAN, RGBColor(0xE0, 0xF2, 0xF1),
     ["Source vs target compare", "Column drift detection", "Type mismatch alerts",
      "Smart recommendations"]),
]

pg_w = Inches(2.45)
pg_h = Inches(1.15)
pg_gap_x = Inches(0.13)
pg_gap_y = Inches(0.12)
pg_x0 = Inches(0.2)
pg_y0 = Inches(0.6)
cols = 5

for idx, (title, clr, light, items) in enumerate(pages):
    r = idx // cols
    c = idx % cols
    px = pg_x0 + c * (pg_w + pg_gap_x)
    py = pg_y0 + r * (pg_h + pg_gap_y)
    box(sl3, px, py, pg_w, Inches(0.26), clr, title, sz=8, fc=WHITE, bold=True)
    mbox(sl3, px, py + Inches(0.28), pg_w, pg_h - Inches(0.28), items,
         sz=6.5, fill=light, line_clr=clr)

# ── Settings + Reconciliation cards (bottom row, 2 extra) ────────────────────
extra_y = pg_y0 + 3 * (pg_h + pg_gap_y)

extra_pages = [
    ("⚙️ Settings", DARK_GRAY, LIGHT_GRAY,
     ["Databricks host + token", "Source DB connection + test", "Metadata catalog/schema",
      "CDC/SDP mode config", "Auto-saves to deployconfig.json"]),
    ("📐 Reconciliation", ORANGE, RGBColor(0xFF, 0xF3, 0xE0),
     ["Source vs Bronze validation", "SUM per numeric column", "PASS / WARN / FAIL status",
      "Dedicated recon catalog", "Failure investigation"]),
    ("📊 deployconfig.json", DEEP_BLUE, RGBColor(0xE3, 0xF2, 0xFD),
     ["databricks_host + token", "source (server, db, user, pass)", "catalogs (bronze, silver, volumes)",
      "metadata_catalog / schema", "cdc_mode + dlt_mode + primary_keys"]),
]

for idx, (title, clr, light, items) in enumerate(extra_pages):
    px = pg_x0 + idx * (pg_w + pg_gap_x)
    box(sl3, px, extra_y, pg_w, Inches(0.26), clr, title, sz=8, fc=WHITE, bold=True)
    mbox(sl3, px, extra_y + Inches(0.28), pg_w, Inches(1.0), items,
         sz=6.5, fill=light, line_clr=clr)

# Pipeline modes summary card
px_mode = pg_x0 + 3 * (pg_w + pg_gap_x)
box(sl3, px_mode, extra_y, pg_w * 2 + pg_gap_x, Inches(0.26), TEAL,
    "Pipeline Execution Modes", sz=8, fc=WHITE, bold=True)
mbox(sl3, px_mode, extra_y + Inches(0.28), pg_w * 2 + pg_gap_x, Inches(1.0),
     ["Standard PySpark Mode (3 stages):",
      "  00_Orchestrator → 01_Extract → 02_Bronze → 03_Silver → 04_Recon → 05_Log",
      "",
      "SDP Mode (2 stages):",
      "  00_Orchestrator → 01_Extract → 02_DLT_Pipeline (Auto Loader + expectations)",
      "",
      "CDC: Change Tracking | Watermark-based incremental | Full load (truncate)"],
     sz=6.5, fill=RGBColor(0xE0, 0xF2, 0xF1), line_clr=TEAL)

# Footer
txt(sl3, Inches(0.2), Inches(7.2), Inches(13), Inches(0.2),
    "SQL → Databricks Migration Studio v2.0  |  14 Feature Pages + Settings + Reconciliation  |  85+ API Endpoints",
    sz=7, fc=MED_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
#  Save
# ═════════════════════════════════════════════════════════════════════════════

output_path = "Reference_Architecture_SQL_to_Databricks_v2.pptx"
prs.save(output_path)
print(f"Saved: {output_path}  (3 slides)")
