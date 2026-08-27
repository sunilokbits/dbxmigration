"""
Generate Update v1 PowerPoint
Documents the v1 Architecture Refactoring:
  - Blueprint decomposition (app.py 2,710 → 75 lines)
  - Static file extraction (index.html 11,548 → 3,790 lines)
  - SQLite persistence layer
  - API versioning (/api/v1/)
  - Project cleanup
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colours ─────────────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_WHITE  = RGBColor(0xFD, 0xFD, 0xFD)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
AZURE_BLUE  = RGBColor(0x00, 0x78, 0xD4)
DEEP_BLUE   = RGBColor(0x0D, 0x47, 0xA1)
DBR_RED     = RGBColor(0xFF, 0x3E, 0x21)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
PURPLE      = RGBColor(0x7B, 0x1F, 0xA2)
TEAL        = RGBColor(0x00, 0x96, 0x88)
ORANGE      = RGBColor(0xE6, 0x51, 0x00)
RED         = RGBColor(0xC6, 0x28, 0x28)
AMBER       = RGBColor(0xFF, 0x8F, 0x00)
INDIGO      = RGBColor(0x30, 0x3F, 0x9F)
INFRA_BLUE  = RGBColor(0x15, 0x65, 0xC0)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────
def box(slide, l, t, w, h, fill, text="", sz=9, fc=WHITE, bold=False,
        align=PP_ALIGN.CENTER, line_clr=None, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_clr:
        s.line.color.rgb = line_clr; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = radius
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(2)
    if text:
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.color.rgb = fc; r.font.bold = bold
    return s


def mbox(slide, l, t, w, h, lines, sz=8, fc=DARK_GRAY, fill=None, line_clr=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_clr:
        s.line.color.rgb = line_clr; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = 0.04
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        r = p.add_run(); r.text = ln; r.font.size = Pt(sz); r.font.color.rgb = fc
    return s


def txt(slide, l, t, w, h, text, sz=9, fc=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(sz); r.font.color.rgb = fc; r.font.bold = bold
    return tb


def arrow_r(slide, x1, y, x2, color=MED_GRAY):
    ln = slide.shapes.add_connector(1, x1, y, x2 - Inches(0.12), y)
    ln.line.color.rgb = color; ln.line.width = Pt(2)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  x2 - Inches(0.14), y - Inches(0.06), Inches(0.14), Inches(0.12))
    tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = color; tri.line.fill.background()


def arrow_d(slide, x, y1, y2, color=MED_GRAY):
    ln = slide.shapes.add_connector(1, x, y1, x, y2 - Inches(0.12))
    ln.line.color.rgb = color; ln.line.width = Pt(1.5)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  x - Inches(0.06), y2 - Inches(0.12), Inches(0.12), Inches(0.12))
    tri.rotation = 180; tri.fill.solid(); tri.fill.fore_color.rgb = color; tri.line.fill.background()


def section_header(slide, x, y, w, title, color):
    box(slide, x, y, w, Inches(0.28), color, title, sz=9, fc=WHITE, bold=True)


def feature_card(slide, x, y, w, h, title, items, color, light):
    box(slide, x, y, w, Inches(0.26), color, title, sz=8, fc=WHITE, bold=True)
    mbox(slide, x, y + Inches(0.28), w, h - Inches(0.28), items,
         sz=6.5, fill=light, line_clr=color)


def stat_badge(slide, x, y, num, label, color):
    box(slide, x, y, Inches(1.4), Inches(0.55), color,
        "", sz=1, fc=WHITE, radius=0.12)
    txt(slide, x + Inches(0.05), y + Inches(0.02), Inches(1.3), Inches(0.3),
        str(num), sz=22, fc=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.05), y + Inches(0.32), Inches(1.3), Inches(0.2),
        label, sz=7, fc=WHITE, bold=False, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — Title Slide
# ═══════════════════════════════════════════════════════════════════════════════

sl1 = prs.slides.add_slide(prs.slide_layouts[6])

# Dark background
bg = sl1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                           Inches(13.333), Inches(7.5))
bg.fill.solid(); bg.fill.fore_color.rgb = RGBColor(0x0F, 0x17, 0x2A)
bg.line.fill.background()

# Accent stripe
box(sl1, Inches(0), Inches(2.8), Inches(13.333), Inches(0.06),
    RGBColor(0x63, 0x66, 0xF1))

# Title
txt(sl1, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.6),
    "Architecture Update v1.0", sz=36, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)

# Subtitle
txt(sl1, Inches(0.8), Inches(2.15), Inches(11.5), Inches(0.5),
    "SQL → Databricks Migration Studio  |  Structural Refactoring & Code Quality",
    sz=16, fc=RGBColor(0xA5, 0xB4, 0xFC), bold=False, align=PP_ALIGN.LEFT)

# Key stats badges
stat_badge(sl1, Inches(0.8),  Inches(3.3), "2,710 → 75", "app.py Lines",    PURPLE)
stat_badge(sl1, Inches(2.5),  Inches(3.3), "11K → 3.8K", "index.html Lines", TEAL)
stat_badge(sl1, Inches(4.2),  Inches(3.3), "14",          "Blueprint Modules", INDIGO)
stat_badge(sl1, Inches(5.9),  Inches(3.3), "110",         "API Routes",       GREEN)
stat_badge(sl1, Inches(7.6),  Inches(3.3), "6",           "Static Files",     ORANGE)
stat_badge(sl1, Inches(9.3),  Inches(3.3), "SQLite",      "Persistence",      INFRA_BLUE)
stat_badge(sl1, Inches(11.0), Inches(3.3), "/api/v1/",    "Versioned API",    AMBER)

# Issues addressed
issues_y = Inches(4.3)
txt(sl1, Inches(0.8), issues_y, Inches(11.5), Inches(0.3),
    "Architecture Review Issues Addressed:", sz=12, fc=RGBColor(0xA5, 0xB4, 0xFC), bold=True,
    align=PP_ALIGN.LEFT)

issue_items = [
    ("#7  Split monolithic app.py into Blueprints (routes/auth.py, routes/convert.py, etc.)",  GREEN),
    ("#8  Split single-file index.html into separate CSS/JS static files",                      GREEN),
    ("#9  Add SQLite persistence layer for in-memory state (MIGRATION_JOBS, _DM_MODELS)",       GREEN),
    ("#10 API versioning — all routes under /api/v1/* with backward-compat redirects",          GREEN),
    ("#11 Project cleanup — removed backups, duplicates, organized test files",                  GREEN),
]
for i, (issue, clr) in enumerate(issue_items):
    iy = issues_y + Inches(0.35) + i * Inches(0.28)
    box(sl1, Inches(0.8), iy, Inches(0.18), Inches(0.18), clr, "✓", sz=8, fc=WHITE, bold=True)
    txt(sl1, Inches(1.1), iy, Inches(11.0), Inches(0.2),
        issue, sz=10, fc=WHITE, align=PP_ALIGN.LEFT)

# Footer
txt(sl1, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.2),
    "April 2026  •  Migration Studio Engineering Team", sz=9, fc=MED_GRAY, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — Before vs After Architecture Diagram
# ═══════════════════════════════════════════════════════════════════════════════

sl2 = prs.slides.add_slide(prs.slide_layouts[6])

# Title bar
box(sl2, Inches(0), Inches(0), Inches(13.333), Inches(0.55), DEEP_BLUE,
    "Architecture Refactoring  —  Before vs After", sz=15, fc=WHITE, bold=True)

# ═════════ LEFT: Before ═════════
before_x = Inches(0.2)
before_w = Inches(6.2)

box(sl2, before_x, Inches(0.7), before_w, Inches(0.3), RED,
    "BEFORE  —  Monolithic Architecture", sz=10, fc=WHITE, bold=True)

# app.py monolith
box(sl2, before_x + Inches(0.3), Inches(1.15), Inches(5.5), Inches(0.3),
    RGBColor(0xFF, 0xEB, 0xEE), "app.py  —  2,710 lines  •  ~100 routes  •  all logic in one file",
    sz=8, fc=RED, bold=True, line_clr=RED)

mono_items = [
    "Authentication + Login",
    "Page serving (index, help, bom)",
    "SP Converter + Upload",
    "Databricks / Unity Catalog",
    "Source DB connection",
    "Data Migration + Jobs",
    "Medallion notebook gen",
    "Self-Healing Bot",
    "Workflow Manager (25+ routes)",
    "Scheduler (CRON jobs)",
    "Reports + Audit + DQ",
    "Schema Compare + Recon",
    "Settings + Deploy Infra",
    "Data Modeling (AI)",
]
row_h = Inches(0.2)
for i, item in enumerate(mono_items):
    ry = Inches(1.5) + i * row_h
    bg = RGBColor(0xFF, 0xEB, 0xEE) if i % 2 == 0 else WHITE
    box(sl2, before_x + Inches(0.5), ry, Inches(5.1), row_h, bg,
        f"  {item}", sz=7, fc=DARK_GRAY, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xEF, 0x9A, 0x9A))

# index.html monolith
box(sl2, before_x + Inches(0.3), Inches(4.4), Inches(5.5), Inches(0.3),
    RGBColor(0xFF, 0xEB, 0xEE), "index.html  —  11,548 lines  •  CSS + JS all inline",
    sz=8, fc=RED, bold=True, line_clr=RED)

html_items = [
    ("Lines 10–1,316",  "1,306 lines of inline CSS"),
    ("Lines 5,050–9,243", "4,193 lines of main JavaScript"),
    ("Lines 9,246–9,298", "52 lines of diagnostic/charts JS"),
    ("Lines 9,304–11,275", "1,971 lines of reports JS"),
    ("Lines 11,309–11,545", "215 lines of supplementary CSS + JS"),
]
for i, (rng, desc) in enumerate(html_items):
    ry = Inches(4.75) + i * Inches(0.22)
    bg = RGBColor(0xFF, 0xEB, 0xEE) if i % 2 == 0 else WHITE
    box(sl2, before_x + Inches(0.5), ry, Inches(1.8), Inches(0.2), bg,
        rng, sz=6.5, fc=MED_GRAY, bold=True, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xEF, 0x9A, 0x9A))
    box(sl2, before_x + Inches(2.35), ry, Inches(3.25), Inches(0.2), bg,
        desc, sz=6.5, fc=DARK_GRAY, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xEF, 0x9A, 0x9A))

# Problems
box(sl2, before_x + Inches(0.3), Inches(5.9), Inches(5.5), Inches(0.24), RED,
    "Problems: Untestable • Merge conflicts • No versioning • State lost on restart",
    sz=7, fc=WHITE, bold=True)

# ═════════ Arrow ═════════
arrow_r(sl2, Inches(6.5), Inches(3.5), Inches(7.0), DEEP_BLUE)
txt(sl2, Inches(6.4), Inches(3.1), Inches(0.8), Inches(0.2),
    "Refactor", sz=8, fc=DEEP_BLUE, bold=True, align=PP_ALIGN.CENTER)

# ═════════ RIGHT: After ═════════
after_x = Inches(7.0)
after_w = Inches(6.1)

box(sl2, after_x, Inches(0.7), after_w, Inches(0.3), GREEN,
    "AFTER  —  Blueprint Architecture", sz=10, fc=WHITE, bold=True)

# app.py (thin)
box(sl2, after_x + Inches(0.2), Inches(1.15), Inches(5.6), Inches(0.3),
    RGBColor(0xE8, 0xF5, 0xE9), "app.py  —  75 lines  •  imports + registers 14 blueprints",
    sz=8, fc=GREEN, bold=True, line_clr=GREEN)

# Blueprint grid (2 columns x 7 rows)
bp_col1 = [
    ("auth.py",       "Login, logout, /api/v1/auth"),
    ("pages.py",      "/, /help, /bom"),
    ("convert.py",    "SP/View/UDF conversion"),
    ("databricks.py", "Test conn, upload, UC"),
    ("source.py",     "Source DB connection"),
    ("migrate.py",    "Data migration + jobs"),
    ("medallion.py",  "Notebook generation"),
]
bp_col2 = [
    ("healer.py",     "Self-healing 15+ routes"),
    ("workflow.py",   "Pipelines, jobs, runs, 25+ routes"),
    ("scheduler.py",  "CRON schedules"),
    ("reports.py",    "Email, audit, DQ"),
    ("schema.py",     "Schema compare + recon"),
    ("settings.py",   "Deploy config + infra"),
    ("datamodel.py",  "AI data modeling"),
]

bp_y_start = Inches(1.55)
for i, (name, desc) in enumerate(bp_col1):
    by = bp_y_start + i * Inches(0.22)
    bg = RGBColor(0xE8, 0xF5, 0xE9) if i % 2 == 0 else WHITE
    box(sl2, after_x + Inches(0.3), by, Inches(1.2), Inches(0.2), bg,
        name, sz=6.5, fc=INDIGO, bold=True, align=PP_ALIGN.LEFT, line_clr=GREEN)
    box(sl2, after_x + Inches(1.55), by, Inches(1.5), Inches(0.2), bg,
        desc, sz=6, fc=DARK_GRAY, align=PP_ALIGN.LEFT, line_clr=GREEN)

for i, (name, desc) in enumerate(bp_col2):
    by = bp_y_start + i * Inches(0.22)
    bg = RGBColor(0xE8, 0xF5, 0xE9) if i % 2 == 0 else WHITE
    box(sl2, after_x + Inches(3.15), by, Inches(1.2), Inches(0.2), bg,
        name, sz=6.5, fc=INDIGO, bold=True, align=PP_ALIGN.LEFT, line_clr=GREEN)
    box(sl2, after_x + Inches(4.4), by, Inches(1.35), Inches(0.2), bg,
        desc, sz=6, fc=DARK_GRAY, align=PP_ALIGN.LEFT, line_clr=GREEN)

# Static files section
static_y = Inches(3.2)
box(sl2, after_x + Inches(0.2), static_y, Inches(5.6), Inches(0.26),
    RGBColor(0xE8, 0xF5, 0xE9), "index.html  →  3,790 lines  +  6 external static files",
    sz=8, fc=GREEN, bold=True, line_clr=GREEN)

static_files = [
    ("main.css",           "1,235 lines", AZURE_BLUE),
    ("main.js",            "3,949 lines", PURPLE),
    ("charts.js",          "49 lines",    TEAL),
    ("reports.js",         "1,806 lines", ORANGE),
    ("supplementary.css",  "12 lines",    INFRA_BLUE),
    ("supplementary.js",   "203 lines",   AMBER),
]
sf_w = Inches(0.88)
sf_gap = Inches(0.05)
for i, (name, lines, clr) in enumerate(static_files):
    sx = after_x + Inches(0.25) + i * (sf_w + sf_gap)
    box(sl2, sx, static_y + Inches(0.3), sf_w, Inches(0.2), clr,
        name, sz=6, fc=WHITE, bold=True)
    box(sl2, sx, static_y + Inches(0.52), sf_w, Inches(0.18),
        RGBColor(0xE8, 0xF5, 0xE9), lines, sz=6, fc=DARK_GRAY, line_clr=clr)

# Persistence + API versioning
pers_y = Inches(4.15)
feature_card(sl2, after_x + Inches(0.2), pers_y, Inches(2.7), Inches(0.78),
    "SQLite Persistence (persistence.py)",
    ["WAL journal mode for concurrency",
     "Thread-local DB connections",
     "Stores MIGRATION_JOBS + _DM_MODELS",
     "init_db() called at app startup"],
    INFRA_BLUE, RGBColor(0xE3, 0xF2, 0xFD))

feature_card(sl2, after_x + Inches(3.05), pers_y, Inches(2.7), Inches(0.78),
    "API Versioning (/api/v1/)",
    ["All 110 routes under /api/v1/ prefix",
     "Backward-compat redirect /api/* → /api/v1/*",
     "HTTP 307 preserves request method",
     "124 frontend fetch() calls updated"],
    AMBER, RGBColor(0xFF, 0xF8, 0xE1))

# Benefits
box(sl2, after_x + Inches(0.2), Inches(5.1), Inches(5.6), Inches(0.24), GREEN,
    "Benefits: Testable • No merge conflicts • Versioned API • Persistent state • Clean separation",
    sz=7, fc=WHITE, bold=True)

# Footer
txt(sl2, Inches(0.2), Inches(7.2), Inches(13), Inches(0.2),
    "Architecture Update v1.0  |  Issues #7–#11  |  app.py 2,710→75 lines  |  index.html 11,548→3,790 lines  |  14 Blueprints  |  110 Routes  |  SQLite Persistence",
    sz=7, fc=MED_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — New Project Structure
# ═══════════════════════════════════════════════════════════════════════════════

sl3 = prs.slides.add_slide(prs.slide_layouts[6])

# Title bar
box(sl3, Inches(0), Inches(0), Inches(13.333), Inches(0.55), TEAL,
    "Clean Project Structure  —  After Refactoring & Cleanup", sz=15, fc=WHITE, bold=True)

# ═════════ COLUMN 1 — File Tree ═════════
col1_x = Inches(0.2)
col1_w = Inches(4.5)

section_header(sl3, col1_x, Inches(0.7), col1_w, "📁 Project File Tree", DARK_GRAY)

tree_lines = [
    "MigrationtoDBXAssetBundle/",
    "├── .gitignore",
    "├── databricks.yml",
    "├── requirements.txt",
    "├── docs/                              # Presentations & docs",
    "├── resources/",
    "│   └── migration_jobs.yml",
    "├── src/",
    "│   ├── notebooks/                     # 9 Databricks notebooks",
    "│   └── sql/init_metadata_tables.sql",
    "└── migration_utility/                 # Flask application",
    "    ├── app.py                         # 75-line orchestrator",
    "    ├── persistence.py                 # SQLite state layer",
    "    ├── __init__.py",
    "    ├── routes/                        # 14 blueprint modules",
    "    │   ├── __init__.py",
    "    │   ├── auth.py        pages.py",
    "    │   ├── convert.py     databricks.py",
    "    │   ├── source.py      migrate.py",
    "    │   ├── medallion.py   healer.py",
    "    │   ├── workflow.py    scheduler.py",
    "    │   ├── reports.py     schema.py",
    "    │   └── settings.py    datamodel.py",
    "    ├── static/                        # 6 extracted files",
    "    │   ├── main.css       main.js",
    "    │   ├── charts.js      reports.js",
    "    │   ├── supplementary.css / .js",
    "    │   └── InsightLogo.bmp",
    "    ├── templates/                     # 4 HTML pages",
    "    │   ├── index.html     login.html",
    "    │   ├── help.html      bom.html",
    "    ├── tests/",
    "    │   └── test_data_modeling.py",
    "    ├── # Core modules:",
    "    ├── sp_converter.py",
    "    ├── stored_procedures.py",
    "    ├── databricks_connector.py",
    "    ├── unity_catalog_executor.py",
    "    ├── data_migrator.py",
    "    ├── data_modeling.py",
    "    ├── medallion_notebooks.py",
    "    ├── metadata_notebooks.py",
    "    ├── self_healing_bot.py",
    "    ├── workflow_manager.py",
    "    ├── AutoInfraCreation.py",
    "    ├── generate_ppt.py",
    "    └── generate_ppt_v2.py",
]

tree_y = Inches(1.05)
for i, line in enumerate(tree_lines):
    ry = tree_y + i * Inches(0.135)
    box(sl3, col1_x + Inches(0.05), ry, col1_w - Inches(0.1), Inches(0.13),
        WHITE, line, sz=5.5, fc=DARK_GRAY, bold=False,
        align=PP_ALIGN.LEFT, line_clr=None)

# ═════════ COLUMN 2 — What Was Removed ═════════
col2_x = Inches(4.9)
col2_w = Inches(3.9)

section_header(sl3, col2_x, Inches(0.7), col2_w, "🗑️ Files Removed (13 items)", RED)

removed_sections = [
    ("Backup Python Files", RED, [
        "app_monolith_backup.py  (2,710 lines)",
        "sp_converter_backup.py",
        "_debug_extract.py  (empty file)",
    ]),
    ("Backup HTML Templates", ORANGE, [
        "index_backup.html",
        "index_pre_redesign.html",
        "index_pre_static_split.html",
    ]),
    ("Duplicate Presentations", AMBER, [
        "DMToolKit - Copy.pptx",
        "Ref_Architecture - Copy.pptx",
        "Sqltodaatbrciks - Copy.pptx",
        "~$Sqltodaatbrciks.pptx  (Office temp)",
    ]),
    ("Runtime Artifacts", MED_GRAY, [
        "schema_compare_results.json",
        "__pycache__/  (31 .pyc files)",
        ".pytest_cache/",
    ]),
]

ry = Inches(1.05)
for title, clr, items in removed_sections:
    box(sl3, col2_x, ry, col2_w, Inches(0.22), clr,
        title, sz=7, fc=WHITE, bold=True)
    ry += Inches(0.24)
    for item in items:
        box(sl3, col2_x, ry, col2_w, Inches(0.18),
            RGBColor(0xFF, 0xEB, 0xEE), f"  ✕  {item}",
            sz=6, fc=DARK_GRAY, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xEF, 0x9A, 0x9A))
        ry += Inches(0.19)
    ry += Inches(0.08)

# ═════════ COLUMN 3 — .gitignore + Stats ═════════
col3_x = Inches(9.05)
col3_w = Inches(4.1)

section_header(sl3, col3_x, Inches(0.7), col3_w, "📋 Updated .gitignore", GREEN)

gitignore_sections = [
    ("Python",     ["__pycache__/", "*.pyc, *.pyo", ".pytest_cache/"]),
    ("Environment",[".venv/", ".env"]),
    ("IDE",        [".idea/", ".vscode/"]),
    ("OS / Office",["~$*", "Thumbs.db", ".DS_Store"]),
    ("Sensitive",  ["deployconfig.json"]),
    ("Runtime",    ["migration_state.db*", "schema_compare_results.json", "job_schedules.json"]),
    ("Databricks", [".databricks/"]),
]

gy = Inches(1.05)
for section, items in gitignore_sections:
    box(sl3, col3_x, gy, Inches(1.2), Inches(0.18), DARK_GRAY,
        section, sz=6, fc=WHITE, bold=True)
    item_text = "  •  ".join(items)
    box(sl3, col3_x + Inches(1.25), gy, Inches(2.8), Inches(0.18),
        RGBColor(0xE8, 0xF5, 0xE9), item_text, sz=5.5, fc=DARK_GRAY,
        align=PP_ALIGN.LEFT, line_clr=GREEN)
    gy += Inches(0.2)

# Summary stats
gy += Inches(0.15)
section_header(sl3, col3_x, gy, col3_w, "📊 Refactoring Summary", INDIGO)
gy += Inches(0.34)

summary_stats = [
    ("app.py",          "2,710 → 75 lines",    "97% reduction",   GREEN),
    ("index.html",      "11,548 → 3,790 lines","67% reduction",   GREEN),
    ("Blueprint files",  "0 → 14 modules",      "routes/*.py",     INDIGO),
    ("Static files",     "0 → 6 files",         "CSS + JS",        TEAL),
    ("Persistence",      "In-memory dicts",      "SQLite WAL",      INFRA_BLUE),
    ("API versioning",   "None (/api/*)",        "/api/v1/*",        AMBER),
    ("Registered routes","~100 (inline)",        "110 (blueprints)", PURPLE),
    ("Test structure",   "Root-level file",      "tests/ directory", ORANGE),
    ("Files removed",    "13 files",             "Backups + dupes",  RED),
]

for i, (metric, before, after, clr) in enumerate(summary_stats):
    sy = gy + i * Inches(0.2)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    box(sl3, col3_x, sy, Inches(1.2), Inches(0.19), bg,
        metric, sz=6, fc=DARK_GRAY, bold=True, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    box(sl3, col3_x + Inches(1.2), sy, Inches(1.4), Inches(0.19), bg,
        before, sz=6, fc=MED_GRAY, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    box(sl3, col3_x + Inches(2.6), sy, Inches(1.5), Inches(0.19), bg,
        after, sz=6, fc=clr, bold=True, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))

# Footer
txt(sl3, Inches(0.2), Inches(7.2), Inches(13), Inches(0.2),
    "Architecture Update v1.0  |  Project Cleanup  |  13 files removed  |  .gitignore hardened  |  tests/ structured",
    sz=7, fc=MED_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
#  Save
# ═══════════════════════════════════════════════════════════════════════════════

output_path = "update_v1.pptx"
prs.save(output_path)
print(f"Saved: {output_path}  (3 slides)")
