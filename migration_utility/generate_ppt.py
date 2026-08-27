"""
Generate a single-slide Reference Architecture PowerPoint for
SQL Server -> Databricks Unity Catalog Migration (Medallion Architecture)
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Colours ─────────────────────────────────────────────────────────────────
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
AZURE_BLUE  = RGBColor(0x00, 0x78, 0xD4)
DBR_RED     = RGBColor(0xFF, 0x3E, 0x21)
LANDING_CLR = RGBColor(0x60, 0x7D, 0x8B)
BRONZE_CLR  = RGBColor(0xCD, 0x7F, 0x32)
SILVER_CLR  = RGBColor(0xAA, 0xA9, 0xAD)
GOLD_CLR    = RGBColor(0xFF, 0xD7, 0x00)
GREEN       = RGBColor(0x2E, 0x7D, 0x32)
PURPLE      = RGBColor(0x7B, 0x1F, 0xA2)
TEAL        = RGBColor(0x00, 0x96, 0x88)
ORANGE      = RGBColor(0xE6, 0x51, 0x00)
RED         = RGBColor(0xC6, 0x28, 0x28)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ── Helpers ─────────────────────────────────────────────────────────────────

def box(slide, l, t, w, h, fill, text="", sz=9, fc=WHITE, bold=False,
        align=PP_ALIGN.CENTER, line_clr=None, radius=0.08):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line_clr:
        s.line.color.rgb = line_clr; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = radius
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(2)
    if text:
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(sz); r.font.color.rgb = fc; r.font.bold = bold
    return s


def mbox(slide, l, t, w, h, lines, sz=8, fc=DARK_GRAY, fill=None, line_clr=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_clr:
        s.line.color.rgb = line_clr; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    s.adjustments[0] = 0.04
    tf = s.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(1)
        r = p.add_run(); r.text = ln; r.font.size = Pt(sz); r.font.color.rgb = fc
    return s


def txt(slide, l, t, w, h, text, sz=9, fc=DARK_GRAY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(sz); r.font.color.rgb = fc; r.font.bold = bold
    return tb


def arrow_r(slide, x1, y, x2, color=MED_GRAY):
    ln = slide.shapes.add_connector(1, x1, y, x2 - Inches(0.12), y)
    ln.line.color.rgb = color; ln.line.width = Pt(2)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  x2 - Inches(0.14), y - Inches(0.06), Inches(0.14), Inches(0.12))
    tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = color; tri.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
#  SINGLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════════

sl = prs.slides.add_slide(prs.slide_layouts[6])

# ── Title Bar ───────────────────────────────────────────────────────────────
box(sl, Inches(0), Inches(0), Inches(13.333), Inches(0.55), AZURE_BLUE,
    "Reference Architecture  --  SQL Server to Databricks Unity Catalog  |  Metadata-Driven Medallion",
    sz=16, fc=WHITE, bold=True)

# ═════════════════════════════════════════════════════════════════════════════
#  ROW 1 — Source -> Orchestrator -> Databricks Medallion Layers
# ═════════════════════════════════════════════════════════════════════════════

row1_y = Inches(0.7)

# ── Source: Azure SQL ────────────────────────────────────────────────────────
src_x = Inches(0.2)
box(sl, src_x, row1_y, Inches(1.6), Inches(0.35), AZURE_BLUE,
    "Azure SQL Server", sz=9, fc=WHITE, bold=True)
mbox(sl, src_x, row1_y + Inches(0.38), Inches(1.6), Inches(0.95),
     ["HR Database", "-------------------", "DimDepartment", "DimEmployee",
      "DimJobRole", "DimLocation", "DimDate"],
     sz=7, fill=RGBColor(0xE3, 0xF2, 0xFD), line_clr=AZURE_BLUE)

# Arrow -> Orchestrator
arrow_r(sl, Inches(1.85), row1_y + Inches(0.55), Inches(2.15), AZURE_BLUE)

# ── Flask AI Workflow Manager ────────────────────────────────────────────────
orch_x = Inches(2.15)
box(sl, orch_x, row1_y, Inches(1.85), Inches(0.35), PURPLE,
    "Flask AI Workflow Mgr", sz=8.5, fc=WHITE, bold=True)
mbox(sl, orch_x, row1_y + Inches(0.38), Inches(1.85), Inches(0.95),
     ["Metadata pipeline config", "SP to PySpark (AI)", "Notebook gen & deploy",
      "Databricks job submit", "Background run polling", "Self-healing recovery"],
     sz=7, fill=RGBColor(0xF3, 0xE5, 0xF5), line_clr=PURPLE)

# Arrow -> Databricks
arrow_r(sl, Inches(4.05), row1_y + Inches(0.55), Inches(4.35), PURPLE)
txt(sl, Inches(4.0), row1_y + Inches(0.2), Inches(0.5), Inches(0.2),
    "REST API", sz=6, fc=MED_GRAY, align=PP_ALIGN.CENTER)

# ── Databricks Container ────────────────────────────────────────────────────
dbr_x = Inches(4.35)
dbr_w = Inches(8.75)
dbr_container = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     dbr_x, row1_y - Inches(0.05), dbr_w, Inches(3.3))
dbr_container.fill.solid()
dbr_container.fill.fore_color.rgb = RGBColor(0xFD, 0xFD, 0xFD)
dbr_container.line.color.rgb = DBR_RED
dbr_container.line.width = Pt(2)
dbr_container.adjustments[0] = 0.02
txt(sl, dbr_x + Inches(0.1), row1_y - Inches(0.02), Inches(3), Inches(0.25),
    "Databricks Unity Catalog", sz=10, fc=DBR_RED, bold=True)

# ── Medallion Layers (inside Databricks) ────────────────────────────────────
layer_y = row1_y + Inches(0.28)
layer_w = Inches(1.85)
gap = Inches(0.2)

layers = [
    ("dev_volumes",   "UC Volumes (Landing)",  LANDING_CLR, RGBColor(0xEC, 0xEF, 0xF1),
     ["/Volumes/dev_volumes/", "  hr/landing/{table}", "Format: Parquet"]),
    ("bronze",        "Bronze (Raw + Audit)",   BRONZE_CLR,  RGBColor(0xFB, 0xE9, 0xD0),
     ["bronze.hr.{table}", "Delta Tables", "+ __src_file, __load_ts", "+ __batch_id"]),
    ("silver",        "Silver (Cleansed)",      SILVER_CLR,  RGBColor(0xEC, 0xEF, 0xF1),
     ["silver.hr.{table}", "Delta Tables", "Dedup + null handling", "Type cast + DQ checks"]),
    ("gold (future)", "Gold (Business KPIs)",   GOLD_CLR,    RGBColor(0xFF, 0xF8, 0xE1),
     ["Aggregated datasets", "BI-ready views", "Power BI / Tableau"]),
]

lx = dbr_x + Inches(0.15)
for i, (cat, subtitle, clr, light, items) in enumerate(layers):
    x = lx + i * (layer_w + gap)
    fc_title = WHITE if cat != "gold (future)" else DARK_GRAY
    box(sl, x, layer_y, layer_w, Inches(0.28), clr, cat, sz=9, fc=fc_title, bold=True)
    box(sl, x, layer_y + Inches(0.3), layer_w, Inches(0.2), light, subtitle,
        sz=7, fc=DARK_GRAY, bold=True, line_clr=clr)
    mbox(sl, x, layer_y + Inches(0.52), layer_w, Inches(0.53), items,
         sz=7, fill=light, line_clr=clr)
    if i < len(layers) - 1:
        arrow_r(sl, x + layer_w, layer_y + Inches(0.6), x + layer_w + gap, clr)

# ── Metadata Tables Row (inside Databricks) ─────────────────────────────────
meta_y = row1_y + Inches(1.45)
box(sl, dbr_x + Inches(0.15), meta_y, dbr_w - Inches(0.3), Inches(0.25),
    GREEN, "Metadata Delta Tables  (Unity Catalog)  +  Reconciliation & Logging Catalogs", sz=8.5, fc=WHITE, bold=True)

meta_tables = ["wf_pipeline_metadata", "wf_job_metadata", "wf_run_history",
               "wf_watermark_metadata", "wf_source_tables",
               "reconciliation.*", "loggingdetails.*"]
mt_w = (dbr_w - Inches(0.3) - Inches(0.1) * 6) / 7
for i, tbl in enumerate(meta_tables):
    mx = dbr_x + Inches(0.15) + i * (mt_w + Inches(0.1))
    clr = GREEN if i < 5 else ORANGE
    box(sl, mx, meta_y + Inches(0.28), mt_w, Inches(0.3),
        RGBColor(0xE8, 0xF5, 0xE9) if i < 5 else RGBColor(0xFF, 0xF3, 0xE0),
        tbl, sz=6, fc=DARK_GRAY, line_clr=clr)

# ── Notebooks Row (inside Databricks) ───────────────────────────────────────
nb_y = meta_y + Inches(0.65)
box(sl, dbr_x + Inches(0.15), nb_y, dbr_w - Inches(0.3), Inches(0.25),
    TEAL, "Auto-Generated Metadata Notebooks  (Standard PySpark + SDP Mode)", sz=8.5, fc=WHITE, bold=True)

notebooks = [
    ("00_Orchestrator",    "Chains all stages\nBreaks on FAILED\nDLT or Spark mode"),
    ("01_Extract",         "JDBC to Azure SQL\nWatermark filtering\nParquet → Volumes"),
    ("02_Bronze / SDP",    "Landing → Bronze\nAudit cols + DQ\nSDP Auto Loader opt."),
    ("03_Silver",          "Bronze → Silver\nDedup + cleanse\n7-check DQ scoring"),
    ("04_Reconciliation",  "Source vs Bronze\nSUM per num col\nPASS/WARN/FAIL"),
    ("05_ExecutionLog",    "Audit trail\nPer-job results\nLogging catalog"),
]
nb_w = (dbr_w - Inches(0.3) - Inches(0.12) * 5) / 6
for i, (name, desc) in enumerate(notebooks):
    nx = dbr_x + Inches(0.15) + i * (nb_w + Inches(0.12))
    box(sl, nx, nb_y + Inches(0.28), nb_w, Inches(0.22),
        RGBColor(0xE0, 0xF2, 0xF1), name, sz=6.5, fc=DARK_GRAY, bold=True, line_clr=TEAL)
    mbox(sl, nx, nb_y + Inches(0.52), nb_w, Inches(0.45), desc.split("\n"),
         sz=6, fill=RGBColor(0xE0, 0xF2, 0xF1), line_clr=TEAL)
    if i < len(notebooks) - 1:
        arrow_r(sl, nx + nb_w, nb_y + Inches(0.55),
                nx + nb_w + Inches(0.12), TEAL)


# ═════════════════════════════════════════════════════════════════════════════
#  ROW 2 — Left: SP Converter + Deploy Infra + Stages | Right: 4 pillars
# ═════════════════════════════════════════════════════════════════════════════

row2_y = Inches(4.15)

# ── Deploy Infrastructure (AutoInfraCreation) ────────────────────────────────
INFRA_BLUE = RGBColor(0x15, 0x65, 0xC0)
box(sl, Inches(0.2), row2_y, Inches(1.95), Inches(0.3), INFRA_BLUE,
    "Deploy Infrastructure", sz=8, fc=WHITE, bold=True)
mbox(sl, Inches(0.2), row2_y + Inches(0.33), Inches(1.95), Inches(0.85),
     ["AutoInfraCreation (Azure SDK)",
      "1. Storage Acct + ADLS Gen2",
      "2. Access Connector + RBAC",
      "3. Storage Credential (UC)",
      "4. External Locations",
      "5. Catalogs + MANAGED LOC",
      "6. UC Volumes (landing)"],
     sz=6.5, fill=RGBColor(0xE3, 0xF2, 0xFD), line_clr=INFRA_BLUE)

# ── AI SP Converter ──────────────────────────────────────────────────────────
box(sl, Inches(0.2), row2_y + Inches(1.3), Inches(0.93), Inches(0.25), ORANGE,
    "AI SP Converter", sz=7, fc=WHITE, bold=True)
mbox(sl, Inches(0.2), row2_y + Inches(1.58), Inches(0.93), Inches(0.35),
     ["OpenAI/Gemini/Claude", "SP → PySpark code"],
     sz=6, fill=RGBColor(0xFF, 0xF3, 0xE0), line_clr=ORANGE)

# ── Pipeline Stages & Rerun ──────────────────────────────────────────────────
box(sl, Inches(1.22), row2_y + Inches(1.3), Inches(0.93), Inches(0.25), PURPLE,
    "Stages & Rerun", sz=7, fc=WHITE, bold=True)
mbox(sl, Inches(1.22), row2_y + Inches(1.58), Inches(0.93), Inches(0.35),
     ["DLT(2)/Spark(3) stages", "Per-job rerun"],
     sz=6, fill=RGBColor(0xF3, 0xE5, 0xF5), line_clr=PURPLE)

# ── 4 Pillars ───────────────────────────────────────────────────────────────
pillars = [
    ("Security", RGBColor(0xC6, 0x28, 0x28), RGBColor(0xFF, 0xEB, 0xEE),
     ["PAT tokens for API auth", "Encrypted JDBC (TLS)",
      "UC RBAC per catalog", "MANAGED LOCATION per catalog"]),
    ("Governance", GREEN, RGBColor(0xE8, 0xF5, 0xE9),
     ["Reconciliation catalog", "Execution logging catalog",
      "Delta versioning", "deployconfig.json driven"]),
    ("Reliability", PURPLE, RGBColor(0xF3, 0xE5, 0xF5),
     ["Per-stage rerun on fail", "Self-healing bot",
      "7-check DQ scoring", "Incremental watermarks"]),
    ("Scalability", TEAL, RGBColor(0xE0, 0xF2, 0xF1),
     ["Metadata-driven config", "SDP + Standard PySpark",
      "Multi-catalog arch", "Extensible to Gold"]),
]

pill_w = Inches(2.4)
pill_gap = Inches(0.15)
pill_x0 = Inches(2.45)

for i, (title, clr, light, items) in enumerate(pillars):
    px = pill_x0 + i * (pill_w + pill_gap)
    box(sl, px, row2_y, pill_w, Inches(0.28), clr, title, sz=8.5, fc=WHITE, bold=True)
    mbox(sl, px, row2_y + Inches(0.3), pill_w, Inches(0.7), items,
         sz=7, fill=light, line_clr=clr)

# ═════════════════════════════════════════════════════════════════════════════
#  ROW 3 — Tech Stack (left) + Data Flow (right-top) + Roadmap (right-bottom)
# ═════════════════════════════════════════════════════════════════════════════

row3_y = Inches(5.25)

# ── Tech Stack (left) ───────────────────────────────────────────────────────
box(sl, Inches(0.2), row3_y, Inches(5.5), Inches(0.25), DARK_GRAY,
    "Technology Stack", sz=9, fc=WHITE, bold=True)

tech_items = [
    ("Azure SQL Server",        "HR source database (5 dimension tables)"),
    ("Python / Flask",          "Web UI: pipeline config, deploy, monitor, rerun"),
    ("Databricks REST API",     "Job submit + polling + notebook deploy"),
    ("UC Volumes + Delta",      "Landing (Parquet) → Bronze/Silver (Delta)"),
    ("Lakeflow Spark Declarative Pipelines",       "Optional SDP mode: Auto Loader + expectations"),
    ("Reconciliation",          "Source vs Bronze aggregate validation per column"),
    ("deployconfig.json",       "Auto-loads catalogs, locations, storage config"),
]

row_h = Inches(0.22)
for i, (tech, desc) in enumerate(tech_items):
    ry = row3_y + Inches(0.28) + i * row_h
    bg = RGBColor(0xE3, 0xF2, 0xFD) if i % 2 == 0 else WHITE
    box(sl, Inches(0.2), ry, Inches(1.9), row_h, bg, tech,
        sz=6.5, fc=DARK_GRAY, bold=True, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    box(sl, Inches(2.1), ry, Inches(3.6), row_h, bg, desc,
        sz=6.5, fc=DARK_GRAY, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))

# ── Data Flow (right-top) ────────────────────────────────────────────────────
flow_x = Inches(6.0)
box(sl, flow_x, row3_y, Inches(7.1), Inches(0.25), AZURE_BLUE,
    "End-to-End Data Flow", sz=9, fc=WHITE, bold=True)

flow_steps = [
    ("Azure\nSQL", AZURE_BLUE),
    ("JDBC\nExtract", PURPLE),
    ("UC Volumes\nLanding", LANDING_CLR),
    ("Bronze\nDelta", BRONZE_CLR),
    ("Recon\nValidation", ORANGE),
    ("Silver\nDelta", SILVER_CLR),
    ("Exec\nLog", GREEN),
]

fs_w = Inches(0.85)
fs_gap = Inches(0.12)
fs_y = row3_y + Inches(0.35)
for i, (label, clr) in enumerate(flow_steps):
    fx = flow_x + Inches(0.15) + i * (fs_w + fs_gap)
    box(sl, fx, fs_y, fs_w, Inches(0.5), clr, label, sz=7, fc=WHITE, bold=True)
    if i < len(flow_steps) - 1:
        arrow_r(sl, fx + fs_w, fs_y + Inches(0.25), fx + fs_w + fs_gap, clr)

# ── Future Roadmap (right-bottom) ────────────────────────────────────────────
road_y = row3_y + Inches(0.95)
box(sl, flow_x, road_y, Inches(7.1), Inches(0.25), DARK_GRAY,
    "Future Roadmap", sz=9, fc=WHITE, bold=True)

roadmap = [
    "Gold Layer: aggregated KPIs for BI dashboards (Power BI / Tableau)",
    "Secret Mgmt: Databricks Secrets / Azure Key Vault (replace PAT tokens)",
    "CI/CD: Azure DevOps / GitHub Actions for notebook deploy & pipeline promotion",
    "Incremental Refresh: CDC-based streaming with Auto Loader / Change Data Feed",
    "Monitoring: Azure Monitor + Databricks SQL Alerts + DQ dashboard",
]

for i, item in enumerate(roadmap):
    ry = road_y + Inches(0.28) + i * Inches(0.2)
    txt(sl, flow_x + Inches(0.15), ry, Inches(6.8), Inches(0.18),
        item, sz=7, fc=DARK_GRAY)

# ── Footer ──────────────────────────────────────────────────────────────────
txt(sl, Inches(0.2), Inches(7.2), Inches(13), Inches(0.2),
    "AI Workflow Manager  |  Metadata-Driven Medallion  |  Reconciliation + Execution Logging  |  SDP + PySpark  |  Per-Stage Rerun  |  deployconfig.json",
    sz=7, fc=MED_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — What's New  (Differences & New Features)
# ═════════════════════════════════════════════════════════════════════════════

sl2 = prs.slides.add_slide(prs.slide_layouts[6])

# ── Title Bar ───────────────────────────────────────────────────────────────
NEW_GREEN = RGBColor(0x2E, 0x7D, 0x32)
box(sl2, Inches(0), Inches(0), Inches(13.333), Inches(0.55), NEW_GREEN,
    "What's New  —  Features Added Since Initial Release",
    sz=16, fc=WHITE, bold=True)

# ── Section helper ──────────────────────────────────────────────────────────
def section_header(slide, x, y, w, title, color):
    box(slide, x, y, w, Inches(0.3), color, title, sz=10, fc=WHITE, bold=True)

def feature_card(slide, x, y, w, h, title, items, color, light):
    box(slide, x, y, w, Inches(0.28), color, title, sz=8.5, fc=WHITE, bold=True)
    mbox(slide, x, y + Inches(0.3), w, h - Inches(0.3), items,
         sz=7, fill=light, line_clr=color)

# ═══════════════ LEFT COLUMN — Major New Features ═══════════════════════════

col1_x = Inches(0.2)
col1_w = Inches(4.2)
cy = Inches(0.7)

section_header(sl2, col1_x, cy, col1_w, "🚀 Major New Features", AZURE_BLUE)
cy += Inches(0.38)

# 1. Deploy Infrastructure
feature_card(sl2, col1_x, cy, col1_w, Inches(1.25),
    "1. Deploy Infrastructure  (AutoInfraCreation)",
    ["Fully automated Azure + Unity Catalog setup via Azure Python SDK",
     "Step 1: ADLS Gen2 Storage Account + Container + Folders (HNS enabled)",
     "Step 2: Access Connector for Databricks (Managed Identity)",
     "Step 3: RBAC — Storage Blob Data Contributor role assignment",
     "Step 4: Unity Catalog Storage Credential",
     "Step 5: External Locations (landing + managed root)",
     "Step 6: Catalogs with MANAGED LOCATION (dev_volumes, bronze, silver)",
     "Step 7: UC Volumes for landing zone (Parquet staging)"],
    RGBColor(0x15, 0x65, 0xC0), RGBColor(0xE3, 0xF2, 0xFD))
cy += Inches(1.32)

# 2. Reconciliation + Execution Log
feature_card(sl2, col1_x, cy, col1_w, Inches(1.0),
    "2. Reconciliation & Execution Logging Notebooks",
    ["04_Meta_Reconciliation — Source vs Bronze aggregate validation",
     "  • Identifies all numeric columns via INFORMATION_SCHEMA",
     "  • Computes SUM per column (JDBC + Delta) → PASS / WARN / FAIL",
     "  • Saves to dedicated reconciliation catalog + Delta table",
     "05_Meta_ExecutionLog — Full audit trail per pipeline run",
     "  • Per-job results saved to loggingdetails.hr.ExecutionLog"],
    ORANGE, RGBColor(0xFF, 0xF3, 0xE0))
cy += Inches(1.07)

# 3. SDP Mode
feature_card(sl2, col1_x, cy, col1_w, Inches(0.85),
    "3. Lakeflow Spark Declarative Pipelines (SDP) Pipeline Mode",
    ["New pipeline_mode: 'dlt' alongside 'standard' PySpark",
     "02_Meta_SDP_Pipeline — Auto Loader + SDP expectations",
     "SDP Orchestrator — Extract → Spark Declarative Pipeline trigger (2 stages)",
     "Standard PySpark remains 3 stages: Extract → Bronze → Silver"],
    DBR_RED, RGBColor(0xFF, 0xEB, 0xEE))
cy += Inches(0.92)

# 4. deployconfig.json
feature_card(sl2, col1_x, cy, col1_w, Inches(0.75),
    "4. deployconfig.json — Centralized Configuration",
    ["Auto-loads catalogs, storage paths, reconciliation & logging config",
     "metadata_notebooks.py reads defaults from deployconfig at import",
     "Eliminates hardcoded values — all infra params in one JSON file"],
    DARK_GRAY, RGBColor(0xF5, 0xF5, 0xF5))

# ═══════════════ RIGHT COLUMN — Pipeline & UI Enhancements ═════════════════

col2_x = Inches(4.6)
col2_w = Inches(4.2)
cy2 = Inches(0.7)

section_header(sl2, col2_x, cy2, col2_w, "🔧 Pipeline & UI Enhancements", PURPLE)
cy2 += Inches(0.38)

# 5. Pipeline Stages
feature_card(sl2, col2_x, cy2, col2_w, Inches(0.85),
    "5. Per-Stage Pipeline Progress Tracking",
    ["Stage progress bar with arrows in Pipeline Studio cards",
     "SDP mode: 2 stages (Extract → SDP Bronze+Silver)",
     "Standard mode: 3 stages (Extract → Bronze → Silver)",
     "Stage-specific status icons: ✅ success / ❌ failed / ⏳ running"],
    PURPLE, RGBColor(0xF3, 0xE5, 0xF5))
cy2 += Inches(0.92)

# 6. Rerun
feature_card(sl2, col2_x, cy2, col2_w, Inches(0.75),
    "6. Per-Job Rerun on Failure",
    ["🔄 Rerun button on each failed stage in Pipeline Studio cards",
     "🔄 Rerun Job button in Pipeline Logs panel per failed run",
     "Reruns from the first failed stage forward (preserves earlier stages)"],
    RED, RGBColor(0xFF, 0xEB, 0xEE))
cy2 += Inches(0.82)

# 7. DQ Framework
feature_card(sl2, col2_x, cy2, col2_w, Inches(0.95),
    "7. Data Quality Framework  (Bronze + Silver)",
    ["Bronze: 5 DQ checks — empty file, null-key, duplicate, schema drift, quarantine",
     "Silver: 7 DQ checks — quarantine filter, null removal, high-null cols,",
     "  dedup, string trim, empty→NULL, row count anomaly detection",
     "DQ Score: % of checks passed → saved to __dq_metrics Delta table",
     "Per-row __dq_status (passed/warn) + __is_quarantined flags"],
    GREEN, RGBColor(0xE8, 0xF5, 0xE9))
cy2 += Inches(1.02)

# 8. MANAGED LOCATION
feature_card(sl2, col2_x, cy2, col2_w, Inches(0.65),
    "8. MANAGED LOCATION for CREATE CATALOG",
    ["Fixed AnalysisException: Metastore storage root URL does not exist",
     "recon_location / log_location threaded through entire stack",
     "Backend → metadata_notebooks → workflow_manager → app → frontend"],
    RGBColor(0xC6, 0x28, 0x28), RGBColor(0xFF, 0xEB, 0xEE))

# ═══════════════ FAR RIGHT — Bug Fixes & Improvements ══════════════════════

col3_x = Inches(9.0)
col3_w = Inches(4.1)
cy3 = Inches(0.7)

section_header(sl2, col3_x, cy3, col3_w, "🐛 Bug Fixes & Improvements", RGBColor(0xC6, 0x28, 0x28))
cy3 += Inches(0.38)

feature_card(sl2, col3_x, cy3, col3_w, Inches(0.75),
    "JDBC ORDER BY Fix",
    ["Removed ORDER BY ORDINAL_POSITION from JDBC subquery",
     "SQL Server forbids ORDER BY in derived tables without TOP/OFFSET",
     "Reconciliation col_query no longer causes SQLServerException"],
    MED_GRAY, RGBColor(0xF5, 0xF5, 0xF5))
cy3 += Inches(0.82)

feature_card(sl2, col3_x, cy3, col3_w, Inches(0.75),
    "Multi-Catalog Validation",
    ["Volumes, Bronze, Silver catalogs must exist before pipeline run",
     "Notebooks run CREATE SCHEMA IF NOT EXISTS (auto-creates schemas)",
     "But CATALOG must already exist — validated in documentation"],
    MED_GRAY, RGBColor(0xF5, 0xF5, 0xF5))
cy3 += Inches(0.82)

feature_card(sl2, col3_x, cy3, col3_w, Inches(0.75),
    "Frontend → Backend Param Fixes",
    ["recon_location / log_location sent in deploy & run-databricks payloads",
     "pipeline_mode sent from UI when creating pipelines (single + bulk)",
     "DLT stage filter added to Job Manager dropdown"],
    MED_GRAY, RGBColor(0xF5, 0xF5, 0xF5))
cy3 += Inches(0.82)

# ── Before vs After summary ─────────────────────────────────────────────────
summary_y = Inches(5.8)
box(sl2, col3_x, summary_y, col3_w, Inches(0.28), AZURE_BLUE,
    "Before vs After", sz=9, fc=WHITE, bold=True)

before_after = [
    ("Notebooks", "4 notebooks", "6 notebooks (+Recon, +ExecLog)"),
    ("Pipeline Modes", "Standard only", "Standard + SDP"),
    ("Stages Visible", "None", "Per-stage progress bars"),
    ("Rerun", "Full rerun only", "Per-job rerun from failure"),
    ("Infra Setup", "Manual", "Automated (AutoInfraCreation)"),
    ("Config", "Hardcoded defaults", "deployconfig.json driven"),
    ("DQ Checks", "Basic", "12 checks + DQ Score + metrics"),
]

ba_y = summary_y + Inches(0.32)
# Header row
box(sl2, col3_x, ba_y, Inches(1.1), Inches(0.2),
    RGBColor(0xE0, 0xE0, 0xE0), "Feature", sz=6.5, fc=DARK_GRAY, bold=True,
    align=PP_ALIGN.LEFT, line_clr=RGBColor(0xBD, 0xBD, 0xBD))
box(sl2, col3_x + Inches(1.1), ba_y, Inches(1.4), Inches(0.2),
    RGBColor(0xFF, 0xEB, 0xEE), "Before", sz=6.5, fc=RED, bold=True,
    align=PP_ALIGN.LEFT, line_clr=RGBColor(0xBD, 0xBD, 0xBD))
box(sl2, col3_x + Inches(2.5), ba_y, Inches(1.6), Inches(0.2),
    RGBColor(0xE8, 0xF5, 0xE9), "After", sz=6.5, fc=GREEN, bold=True,
    align=PP_ALIGN.LEFT, line_clr=RGBColor(0xBD, 0xBD, 0xBD))
ba_y += Inches(0.2)

for feat, before, after in before_after:
    bg = WHITE if before_after.index((feat, before, after)) % 2 == 0 else RGBColor(0xFA, 0xFA, 0xFA)
    box(sl2, col3_x, ba_y, Inches(1.1), Inches(0.18), bg, feat,
        sz=6, fc=DARK_GRAY, bold=True, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    box(sl2, col3_x + Inches(1.1), ba_y, Inches(1.4), Inches(0.18), bg, before,
        sz=6, fc=MED_GRAY, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    box(sl2, col3_x + Inches(2.5), ba_y, Inches(1.6), Inches(0.18), bg, after,
        sz=6, fc=GREEN, align=PP_ALIGN.LEFT, line_clr=RGBColor(0xE0, 0xE0, 0xE0))
    ba_y += Inches(0.18)

# ── Footer ──────────────────────────────────────────────────────────────────
txt(sl2, Inches(0.2), Inches(7.2), Inches(13), Inches(0.2),
    "SQL to Databricks Migration  |  What's New — All changes since initial architecture release",
    sz=7, fc=MED_GRAY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
#  Save
# ═════════════════════════════════════════════════════════════════════════════

output_path = "Reference_Architecture_SQL_to_Databricks.pptx"
prs.save(output_path)
print(f"Saved: {output_path}  (2 slides)")
